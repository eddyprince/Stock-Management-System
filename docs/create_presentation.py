"""
Generate PowerPoint: Stock Management Control System
- System description
- Use case diagram (summary)
- Class diagram (summary)
- Data flow diagram (summary)
- Sequence diagram (summary)

Run: pip install python-pptx && python create_presentation.py
Output: StockManagementSystem_Description.pptx (in docs/ or project root)
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Install: pip install python-pptx")
    raise

# Theme colors (teal/slate) — python-pptx 1.x uses RGBColor
PRIMARY = RGBColor(20, 184, 166)   # primary-600
DARK = RGBColor(15, 23, 42)        # slate-900

def add_title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    left = Inches(0.5)
    top = Inches(1.5)
    w = prs.slide_width - Inches(1)
    tf = slide.shapes.add_textbox(left, top, w, Inches(1.2))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    if subtitle:
        tf2 = slide.shapes.add_textbox(left, Inches(2.8), w, Inches(0.8))
        tf2.text_frame.paragraphs[0].text = subtitle
        tf2.text_frame.paragraphs[0].font.size = Pt(18)
        tf2.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)
    return slide

def add_content_slide(prs, title, bullets):
    # Use blank layout (index 6 or 5 depending on theme)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    left = Inches(0.5)
    top = Inches(0.4)
    w = prs.slide_width - Inches(1)
    # Slide title
    tf = slide.shapes.add_textbox(left, top, w, Inches(0.7))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    # Content: flatten (title, [subs]) into main + sub bullets
    flat = []
    for item in bullets:
        if isinstance(item, tuple) and len(item) >= 2:
            flat.append(item[0])
            for sub in item[1]:
                flat.append("  • " + sub)
        else:
            flat.append(item if isinstance(item, str) else item[0])
    body = slide.shapes.add_textbox(left, Inches(1.1), w, Inches(6))
    text_frame = body.text_frame
    text_frame.word_wrap = True
    for i, line in enumerate(flat):
        para = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
        para.text = line
        para.font.size = Pt(15)
        para.space_after = Pt(4)
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title slide
    add_title_slide(prs,
        "Stock Management Control System",
        "System description, use case, class, data flow & sequence diagrams"
    )

    # 2. System description
    add_content_slide(prs, "System Description", [
        "Purpose: Web app for authorized users to view and manage stock (expired, in stock, out of stock, damaged).",
        "Users log in by role: Account, Stock Manager, Admin, Director.",
        "High-level flow:",
        ("", [
            "User visits app → SPA loads (Vue).",
            "Login → backend validates, returns JWT; frontend stores token.",
            "Role-based access: router and API enforce permissions.",
            "Stock data: products with supplier, total supplied, expiry.",
            "Reports: Director/Admin view dashboards (expired, in/out stock, damaged).",
        ]),
        "Tech: Vue 3, Vite, Pinia, Vue Router (frontend); Node.js, Express, MongoDB (backend).",
    ])

    # 3. Use case diagram
    add_content_slide(prs, "Use Case Diagram", [
        "Actors: User, Admin, Stock Manager, Director, Account.",
        "Main use cases:",
        ("", [
            "Login / Logout — all roles.",
            "View Stock Info, View Expired Stock, View Damaged Qty — Director, Stock Manager, Admin.",
            "Manage Products — Stock Manager, Admin.",
            "Record Supply / Supplier — Stock Manager, Admin (where from, total supplied).",
            "View Reports — Director, Admin.",
            "Manage Users — Admin only.",
            "Account Operations — Account role.",
        ]),
        "Visual diagram: see docs/DIAGRAMS.md (Mermaid) or docs/SYSTEM_OVERVIEW.md.",
    ])

    # 4. Class diagram
    add_content_slide(prs, "Class Diagram", [
        "User: _id, username, passwordHash, role, email, createdAt.",
        "Product: _id, name, sku, quantityInStock, quantityDamaged, totalSupplied, supplierId, expiryDate; isExpired().",
        "Supplier: _id, name, contact.",
        "StockTransaction: _id, productId, type, quantity, amount, userId, createdAt.",
        "Relationships:",
        ("", [
            "Product → Supplier (from).",
            "Product → StockTransaction (has).",
            "User → StockTransaction (creates).",
        ]),
        "Visual diagram: see docs/DIAGRAMS.md (Mermaid).",
    ])

    # 5. Data flow diagram
    add_content_slide(prs, "Data Flow Diagram", [
        "Level 0 (context): User ↔ Stock Management System ↔ MongoDB.",
        "Level 1 (main processes):",
        ("", [
            "User → Login Form → Auth API → MongoDB.",
            "User → Dashboard/Stock → Products/Stock API → MongoDB.",
            "User → Product/Supply Forms → Products/Stock API → MongoDB.",
            "User → Reports View → Reports API → MongoDB.",
            "Data: Login/CRUD/Reports to System; System reads/writes DB; responses/UI to User.",
        ]),
        "Visual diagram: see docs/DIAGRAMS.md (Mermaid).",
    ])

    # 6. Sequence diagram
    add_content_slide(prs, "Sequence Diagram (Login & View Stock)", [
        "1. User opens app → Vue loads.",
        "2. Router checks route; no token → redirect to Login.",
        "3. User enters credentials → Vue sends POST /auth/login to API.",
        "4. API verifies user in MongoDB → returns JWT.",
        "5. Vue stores token, redirects (e.g. to Dashboard).",
        "6. Vue sends GET /products with JWT.",
        "7. API queries MongoDB → returns stock data.",
        "8. Vue displays dashboard.",
        "Visual diagram: see docs/DIAGRAMS.md (Mermaid).",
    ])

    # 7. Reference
    add_content_slide(prs, "Reference", [
        "All diagram source code (Mermaid) and full system overview:",
        "  • docs/SYSTEM_OVERVIEW.md",
        "  • docs/DIAGRAMS.md",
        "To render Mermaid diagrams as images: use GitHub, Mermaid Live Editor, or VS Code Mermaid extension.",
    ])

    out = Path(__file__).parent / "StockManagementSystem_Description.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
